# parse_pptx.py —— PowerPoint（.pptx）解析（python-pptx）。
#
# 能力：
#   - 文本提取（按幻灯片分组，标题为 "## 幻灯片 N"）
#   - 图片提取（shape.image）与视频/音频提取（shape.media_format）
#   - 公式检测：PPT 内的 OMML 公式 python-pptx 拿不到字符，检测到则 warning 提示
#
# usage: python3 parse_pptx.py <input_abs> <out_json_abs> <media_dir_abs>
import sys

import common


def main() -> int:
    if len(sys.argv) != 4:
        common.usage("parse_pptx.py")
    input_path, out_path, media_dir = sys.argv[1], sys.argv[2], sys.argv[3]
    media_ref = common.media_rel_prefix(media_dir)

    warnings: list[str] = []
    media: list[dict] = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
    except ImportError:
        common.emit(out_path, "", "", [], False,
                    ["沙盒缺少 python-pptx（镜像构建期未安装），无法解析 PPT"])
        return 0

    try:
        prs = Presentation(input_path)
    except Exception as e:  # noqa: BLE001
        common.emit(out_path, "", "", [], False, ["PPT 打开失败: %s" % e])
        return 0

    md: list[str] = []
    title = ""
    for idx, slide in enumerate(prs.slides):
        slide_lines: list[str] = []
        # 公式检测：幻灯片 XML 含 m:oMath（Office Math）→ 文本层拿不到公式。
        try:
            if slide.part.element.find(".//" + qn("m:oMath")) is not None:
                warnings.append("检测到 PPT 含公式（第 %d 页），公式字符无法从文本层还原，建议对照原课件" % (idx + 1))
        except Exception:  # noqa: BLE001
            pass

        img_n = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_lines.append(text)
            try:
                stype = shape.shape_type
            except Exception:  # noqa: BLE001
                stype = None

            if stype == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    ext = (img.ext or "png").lower()
                    ref = common.save_bytes(
                        img.blob, media_dir, "slide%d_img%d.%s" % (idx + 1, img_n + 1, ext), media_ref)
                    media.append({"type": "image", "path": ref, "alt": "幻灯片 %d 图 %d" % (idx + 1, img_n + 1)})
                    slide_lines.append("![幻灯片 %d 图 %d](%s)" % (idx + 1, img_n + 1, ref))
                    img_n += 1
                except Exception:  # noqa: BLE001
                    continue
            elif stype == MSO_SHAPE_TYPE.MEDIA:
                try:
                    mf = shape.media_format
                    blob = mf.blob
                    ext = (mf.ext or "bin").lower()
                    ref = common.save_bytes(
                        blob, media_dir, "slide%d_media%d.%s" % (idx + 1, img_n + 1, ext), media_ref)
                    mtype = common.media_type(ext)
                    media.append({"type": mtype, "path": ref, "alt": ""})
                    warnings.append("第 %d 页含媒体附件（%s），已提取保存，内容未解析" % (idx + 1, ref))
                    img_n += 1
                except Exception:  # python-pptx 版本不支持 media_format
                    warnings.append("第 %d 页含内嵌媒体，但沙盒 python-pptx 版本不支持提取" % (idx + 1))

        if idx == 0:
            # 首页非空文本尝试作为标题。
            for line in slide_lines:
                if line:
                    title = line[:120]
                    break
        md.append("## 幻灯片 %d" % (idx + 1))
        md.extend(slide_lines)

    common.emit(out_path, title, "\n\n".join(md), media, False, warnings)
    return 0


if __name__ == "__main__":
    sys.exit(main())
