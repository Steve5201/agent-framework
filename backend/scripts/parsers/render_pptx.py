# render_pptx.py —— 按 DocumentSpec 渲染演示文稿（.pptx）（P4-D / P4-J 阶段2）。
#
# 技术栈：python-pptx（沙盒镜像已预装 python-pptx==1.0.2，见 deploy/Dockerfile）。
# 能力：
#   - 16:9 版式：首页 = 标题 + 副标题；每节 = 一页（标题 + 正文）；
#   - 富文本块（P4-J 阶段2）：文本类块（paragraph/list/code 等）合并进正文
#     文本框；image/formula 块转 PNG 后插入；table 块用原生表格；
#   - blocks 为空时回退旧契约 body 轻量标记（- 列表 / # ## 加粗段 / 段落）；
#   - 页码页脚（python-pptx 无原生页脚，用每页底部小文本框实现）。
#
# usage: python3 render_pptx.py <spec_json_abs> <out_pptx_abs>
import os
import sys

import render_common


def _write_body(tf, body: str) -> None:
    """把 body 按轻量标记写入文本框（旧契约回退：首段即用，后续逐行 add_paragraph）。"""
    from pptx.util import Pt

    first = True
    for kind, text in render_common.parse_body(body):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.font.size = Pt(20)
        if kind == "bullet":
            p.level = 0
            # python-pptx 的 placeholder 默认带项目符号；手动文本框则显式加前缀。
            p.text = "• " + text
        elif kind == "h2":
            p.level = 0
            p.font.size = Pt(24)
            p.font.bold = True
        elif kind == "h3":
            p.level = 1
            p.font.size = Pt(22)
            p.font.bold = True
        else:
            p.level = 0


def _render_blocks_pptx(prs, slide, blocks) -> None:
    """按 blocks 渲染一节内容：文本类块合并进正文框，图片/公式/表格依次堆叠。"""
    from pptx.util import Emu, Inches, Pt

    slide_w = prs.slide_width
    x_left = Inches(0.3)
    body_top = Inches(1.1)
    body_bottom = Inches(6.9)
    body_width = slide_w - Inches(0.6)

    text_blocks = []
    media_blocks = []
    for b in blocks:
        kind = b.get("type", "")
        if kind in ("paragraph", "p", "h2", "h3", "bullet", "list", "code"):
            text_blocks.append(b)
        else:
            media_blocks.append(b)

    # ---- 文本类块：合并进一个正文文本框 ----
    y = body_top
    if text_blocks:
        tb = slide.shapes.add_textbox(x_left, body_top, body_width, body_bottom - body_top)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for b in text_blocks:
            kind = b.get("type", "")
            if kind == "list":
                items = b.get("items") or []
            else:
                items = [b.get("text", "")]
            for item in items:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = item
                p.font.size = Pt(20)
                if kind in ("h2",):
                    p.font.size = Pt(24)
                    p.font.bold = True
                elif kind in ("h3",):
                    p.font.size = Pt(22)
                    p.font.bold = True
                elif kind in ("bullet", "list"):
                    p.text = "• " + item
                elif kind == "code":
                    p.font.name = "Consolas"
                    p.font.size = Pt(14)
        y = body_top + Inches(0.4)

    # ---- 非文本块（图片/公式/表格）：从文本区下方垂直堆叠 ----
    work_dir = os.getcwd()
    for b in media_blocks:
        kind = b.get("type", "")
        if kind == "image":
            try:
                png, _, _ = render_common.asset_png(b.get("src", ""), work_dir)
            except ValueError:
                continue
            w_px = b.get("width", 0)
            width_in = Inches(4) if w_px <= 0 else Inches(min(w_px / 96.0, 6))
            pic = slide.shapes.add_picture(png, x_left, y, width=width_in)
            # 居中：按实际落点宽度修正 x。
            pic.left = Emu(int((slide_w - pic.width) / 2))
            y += pic.height + Inches(0.2)
            if b.get("caption"):
                cb = slide.shapes.add_textbox(x_left, y, body_width, Inches(0.35))
                cap_tf = cb.text_frame
                cap_tf.word_wrap = True
                cap_tf.text = b["caption"]
                cap_tf.paragraphs[0].font.size = Pt(12)
                cap_tf.paragraphs[0].alignment = 2  # PP_ALIGN.CENTER
                y += Inches(0.4)
        elif kind == "formula":
            # render=native → OMML 原生公式（PPT 2010+ a14:m）；结构不支持回退图片。
            if b.get("render", "image") == "native":
                try:
                    tb = slide.shapes.add_textbox(x_left, y, body_width, Inches(0.6))
                    tf = tb.text_frame
                    tf.word_wrap = False
                    p = tf.paragraphs[0]
                    p.alignment = 2  # PP_ALIGN.CENTER
                    render_common.omml_into_pptx_paragraph(
                        p, render_common.latex_to_omml(b.get("text", ""))
                    )
                    y += Inches(0.7)
                    continue
                except Exception:
                    pass  # OMML 注入失败 → 回退图片
            try:
                png = render_common.formula_png(b.get("text", ""), work_dir)
            except ValueError:
                continue
            pic = slide.shapes.add_picture(png, x_left, y, height=Inches(0.8))
            pic.left = Emu(int((slide_w - pic.width) / 2))
            y += pic.height + Inches(0.2)
        elif kind == "table":
            headers = b.get("headers") or []
            rows = b.get("rows") or []
            nrows, ncols = len(rows) + 1, len(headers)
            if nrows > 0 and ncols > 0:
                tbl_shape = slide.shapes.add_table(
                    nrows, ncols, x_left, y, body_width, Inches(max(0.4 * nrows, 0.5))
                )
                table = tbl_shape.table
                for j, h in enumerate(headers):
                    table.cell(0, j).text = h
                for i, row in enumerate(rows, start=1):
                    for j, val in enumerate(row):
                        table.cell(i, j).text = val
                y += tbl_shape.height + Inches(0.2)


def _add_cover(prs, slide, spec) -> None:
    """封面：标题/副标题居中（Blank 版式无占位符，手动文本框，居中布局）。"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.4), prs.slide_width - Inches(1.0), Inches(1.4)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = spec["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    if spec["subtitle"]:
        tb2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.0), prs.slide_width - Inches(1.0), Inches(0.8)
        )
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.text = spec["subtitle"]
        p2.font.size = Pt(22)


def _add_section_title(prs, slide, heading: str) -> None:
    """章节标题：手动文本框（Blank 版式无占位符）。"""
    from pptx.util import Inches, Pt

    tb = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.3), prs.slide_width - Inches(0.6), Inches(0.8)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(30)
    p.font.bold = True


def main() -> int:
    if len(sys.argv) != 3:
        render_common.usage("render_pptx.py")
    spec_path, out_path = sys.argv[1], sys.argv[2]
    try:
        spec = render_common.load_spec(spec_path)
    except ValueError as e:
        render_common.emit_error("render_pptx: %s" % e)
        return 1

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        render_common.emit_error("render_pptx: 沙盒缺少 python-pptx（%s）" % e)
        return 1

    prs = Presentation()
    # 16:9 版式（13.333 x 7.5 英寸）。
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    # Blank 版式：全空页，不带模板占位符/页码等自带元素（P4-L 修复），
    # 所有内容由脚本手动文本框渲染，保证观感可控。
    blank = prs.slide_layouts[6]

    x_left = Inches(0.3)
    body_width = prs.slide_width - Inches(0.6)

    # 封面：标题 + 副标题（居中）。
    slide = prs.slides.add_slide(blank)
    _add_cover(prs, slide, spec)

    # 每节一页：blocks 非空优先（全块走 _render_blocks_pptx，含纯文本块），
    # blocks 为空回退 body 轻量标记。
    for sec in spec["sections"]:
        slide = prs.slides.add_slide(blank)
        _add_section_title(prs, slide, sec["heading"])
        if sec.get("blocks"):
            _render_blocks_pptx(prs, slide, sec["blocks"])
        else:
            body_tb = slide.shapes.add_textbox(
                x_left, Inches(1.35), body_width, Inches(5.6)
            )
            _write_body(body_tb.text_frame, sec["body"])

    try:
        prs.save(out_path)
    except OSError as e:
        render_common.emit_error("render_pptx: 保存失败（%s）" % e)
        return 1

    size = render_common.file_size(out_path)
    if size <= 0:
        render_common.emit_error("render_pptx: 产物为空或不可读")
        return 1
    render_common.emit_ok(out_path, size)
    return 0


if __name__ == "__main__":
    sys.exit(main())
